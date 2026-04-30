from __future__ import annotations

import base64
import json
import tempfile
import unittest
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from image2ppt_exact import (
    BlueprintRebuildConfig,
    EditableConfig,
    ExportConfig,
    FullRebuildConfig,
    ImageSvgEditableConfig,
    SvgNativeRebuildConfig,
    create_editable_text_pptx,
    export_exact_deck,
    rebuild_from_blueprint,
    run_full_rebuild_pipeline,
    run_image_svg_editable_pipeline,
    rebuild_from_svg_native,
)
from image2ppt_exact.cli import build_parser
from image2ppt_exact.editable import (
    collect_editable_blocks,
    load_ocr_locks,
    rapid_result_to_blocks,
)


TINY_PNG = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8"
    "/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
)


class ExporterTests(unittest.TestCase):
    def test_exports_svg_preview_and_log(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            src = root / "slides"
            out = root / "out"
            src.mkdir()
            (src / "slide_01.png").write_bytes(base64.b64decode(TINY_PNG))

            result = export_exact_deck(
                ExportConfig(src=src, out=out, create_pptx=False)
            )

            self.assertEqual(result.slide_count, 1)
            self.assertTrue((out / "slides_svg" / "slide_01.svg").exists())
            self.assertTrue((out / "index.html").exists())
            self.assertTrue((out / "run-log.md").exists())
            svg = (out / "slides_svg" / "slide_01.svg").read_text(
                encoding="utf-8"
            )
            self.assertIn("data:image/png;base64", svg)

    def test_creates_editable_text_pptx_from_json(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            src = root / "slides"
            ocr = root / "ocr"
            src.mkdir()
            ocr.mkdir()
            (src / "slide_01.png").write_bytes(base64.b64decode(TINY_PNG))
            (ocr / "slide_01.json").write_text(
                '{"blocks":[{"text":"Hello","bbox":[0,0,100,40]}]}',
                encoding="utf-8",
            )
            pptx = root / "editable.pptx"

            create_editable_text_pptx(
                EditableConfig(
                    src=src,
                    ocr_dir=ocr,
                    pptx_path=pptx,
                    width=1920,
                    height=1080,
                    background="blank",
                )
            )

            self.assertTrue(pptx.exists())

    def test_editable_routes_default_to_blank_background(self) -> None:
        parser = build_parser()

        editable_args = parser.parse_args(
            [
                "editable",
                "slides",
                "--ocr",
                "ocr",
                "--pptx",
                "editable.pptx",
            ]
        )
        pipeline_args = parser.parse_args(
            ["image-svg-editable", "slides", "--out", "out"]
        )
        full_args = parser.parse_args(["full-rebuild", "slides", "--out", "out"])

        self.assertEqual(EditableConfig(Path("s"), Path("o"), Path("p")).background, "blank")
        self.assertEqual(ImageSvgEditableConfig(Path("s"), Path("out")).background, "blank")
        self.assertEqual(FullRebuildConfig(Path("s"), Path("out")).background, "blank")
        self.assertEqual(editable_args.background, "blank")
        self.assertEqual(pipeline_args.background, "blank")
        self.assertEqual(full_args.background, "blank")

    def test_rapidocr_result_to_blocks(self) -> None:
        blocks = rapid_result_to_blocks(
            [
                (
                    [[1, 2], [11, 2], [11, 8], [1, 8]],
                    "Hello",
                    "0.91",
                )
            ]
        )

        self.assertEqual(blocks[0]["text"], "Hello")
        self.assertEqual(blocks[0]["bbox"], [1.0, 2.0, 10.0, 6.0])
        self.assertAlmostEqual(blocks[0]["confidence"], 0.91)

    def test_default_editable_output_does_not_keep_text_bearing_image_background(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            src = root / "slides"
            ocr = root / "ocr"
            src.mkdir()
            ocr.mkdir()
            (src / "slide_01.png").write_bytes(base64.b64decode(TINY_PNG))
            (ocr / "slide_01.json").write_text(
                '{"blocks":[{"text":"Hello","bbox":[0,0,100,40]}]}',
                encoding="utf-8",
            )
            pptx = root / "editable.pptx"

            create_editable_text_pptx(
                EditableConfig(src=src, ocr_dir=ocr, pptx_path=pptx)
            )

            prs = Presentation(str(pptx))
            pictures = [
                shape
                for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ]
            text_boxes = [
                shape
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ]
            self.assertEqual(len(pictures), 0)
            self.assertEqual(len(text_boxes), 1)

    def test_redacted_background_removes_ocr_text_regions_before_overlay(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            src = root / "slides"
            ocr = root / "ocr"
            src.mkdir()
            ocr.mkdir()
            img = Image.new("RGB", (1920, 1080), "white")
            draw = ImageDraw.Draw(img)
            draw.rectangle((10, 10, 70, 34), fill="black")
            img.save(src / "slide_01.png")
            (ocr / "slide_01.json").write_text(
                '{"blocks":[{"text":"Hello","bbox":[10,10,60,24]}]}',
                encoding="utf-8",
            )
            pptx = root / "editable-redacted.pptx"

            create_editable_text_pptx(
                EditableConfig(
                    src=src,
                    ocr_dir=ocr,
                    pptx_path=pptx,
                    width=1920,
                    height=1080,
                    background="redact",
                )
            )

            prs = Presentation(str(pptx))
            pictures = [
                shape
                for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ]
            text_boxes = [
                shape
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ]
            self.assertEqual(len(pictures), 1)
            self.assertEqual(len(text_boxes), 1)

            redacted_bg = root / "editable-redacted_assets" / "slide_01_redacted.png"
            self.assertTrue(redacted_bg.exists())
            with Image.open(redacted_bg) as redacted:
                # Center of the original black text block should be cleared to the
                # surrounding white background before editable text is overlaid.
                self.assertEqual(redacted.getpixel((20, 20)), (255, 255, 255))

    def test_filters_small_ocr_blocks_before_textbox_and_redaction(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            src = root / "slides"
            ocr = root / "ocr"
            src.mkdir()
            ocr.mkdir()
            img = Image.new("RGB", (400, 200), "white")
            draw = ImageDraw.Draw(img)
            draw.rectangle((10, 10, 40, 18), fill="black")
            draw.rectangle((70, 20, 170, 60), fill="black")
            img.save(src / "slide_01.png")
            (ocr / "slide_01.json").write_text(
                json.dumps(
                    {
                        "blocks": [
                            {"text": "Tiny", "bbox": [10, 10, 30, 8]},
                            {"text": "Main", "bbox": [70, 20, 100, 40]},
                        ]
                    }
                ),
                encoding="utf-8",
            )
            pptx = root / "editable-redacted.pptx"

            create_editable_text_pptx(
                EditableConfig(
                    src=src,
                    ocr_dir=ocr,
                    pptx_path=pptx,
                    width=400,
                    height=200,
                    background="redact",
                    min_text_height=12,
                    min_text_area=500,
                )
            )

            prs = Presentation(str(pptx))
            texts = [
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ]
            self.assertEqual(texts, ["Main"])

            redacted_bg = root / "editable-redacted_assets" / "slide_01_redacted.png"
            with Image.open(redacted_bg) as redacted:
                self.assertEqual(redacted.getpixel((20, 14)), (0, 0, 0))
                self.assertEqual(redacted.getpixel((90, 40)), (255, 255, 255))

    def test_lock_file_skips_conversion_and_redaction_inside_locked_regions(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            src = root / "slides"
            ocr = root / "ocr"
            src.mkdir()
            ocr.mkdir()
            img = Image.new("RGB", (400, 200), "white")
            draw = ImageDraw.Draw(img)
            draw.rectangle((20, 20, 80, 50), fill="black")
            draw.rectangle((120, 20, 180, 50), fill="black")
            img.save(src / "slide_01.png")
            (ocr / "slide_01.json").write_text(
                json.dumps(
                    {
                        "blocks": [
                            {"text": "Locked", "bbox": [20, 20, 60, 30]},
                            {"text": "Editable", "bbox": [120, 20, 60, 30]},
                        ]
                    }
                ),
                encoding="utf-8",
            )
            lock_file = root / "ocr-locks.json"
            lock_file.write_text(
                json.dumps(
                    {
                        "slides": {
                            "slide_01": [
                                {"x": 0, "y": 0, "w": 100, "h": 100}
                            ]
                        }
                    }
                ),
                encoding="utf-8",
            )
            pptx = root / "editable-redacted.pptx"

            create_editable_text_pptx(
                EditableConfig(
                    src=src,
                    ocr_dir=ocr,
                    pptx_path=pptx,
                    width=400,
                    height=200,
                    background="redact",
                    lock_file=lock_file,
                )
            )

            prs = Presentation(str(pptx))
            texts = [
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ]
            self.assertEqual(texts, ["Editable"])

            redacted_bg = root / "editable-redacted_assets" / "slide_01_redacted.png"
            with Image.open(redacted_bg) as redacted:
                self.assertEqual(redacted.getpixel((40, 35)), (0, 0, 0))
                self.assertEqual(redacted.getpixel((140, 35)), (255, 255, 255))

    def test_collect_editable_blocks_reports_skip_reasons(self) -> None:
        blocks = [
            {"text": "Tiny", "bbox": [0, 0, 20, 8]},
            {"text": "SmallArea", "bbox": [30, 0, 20, 20]},
            {"text": "Locked", "bbox": [70, 0, 40, 20]},
            {"text": "Main", "bbox": [120, 0, 80, 30]},
        ]
        locks = {"slide_01": [{"x": 60, "y": 0, "w": 60, "h": 60}]}

        editable, stats = collect_editable_blocks(
            blocks,
            slide_stem="slide_01",
            width=240,
            height=120,
            min_text_height=10,
            min_text_area=500,
            locks=locks,
        )

        self.assertEqual([block["text"] for block in editable], ["Main"])
        self.assertEqual(stats.total_blocks, 4)
        self.assertEqual(stats.editable_blocks, 1)
        self.assertEqual(stats.skipped_by_height, 1)
        self.assertEqual(stats.skipped_by_area, 1)
        self.assertEqual(stats.skipped_by_lock, 1)

    def test_load_ocr_locks_accepts_slide_map_and_global_regions(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            lock_file = Path(tmp) / "ocr-locks.json"
            lock_file.write_text(
                json.dumps(
                    {
                        "regions": [{"x": 1, "y": 2, "w": 3, "h": 4}],
                        "slides": {
                            "slide_02": [{"left": 5, "top": 6, "width": 7, "height": 8}]
                        },
                    }
                ),
                encoding="utf-8",
            )

            locks = load_ocr_locks(lock_file)

            self.assertEqual(locks["*"], [{"x": 1.0, "y": 2.0, "w": 3.0, "h": 4.0}])
            self.assertEqual(
                locks["slide_02"],
                [{"x": 5.0, "y": 6.0, "w": 7.0, "h": 8.0}],
            )

    def test_image_svg_editable_pipeline_with_existing_ocr(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            src = root / "slides"
            ocr = root / "ocr"
            out = root / "out"
            src.mkdir()
            ocr.mkdir()
            (src / "slide_01.png").write_bytes(base64.b64decode(TINY_PNG))
            (ocr / "slide_01.json").write_text(
                '{"blocks":[{"text":"Hello","bbox":[0,0,100,40]}]}',
                encoding="utf-8",
            )

            result = run_image_svg_editable_pipeline(
                ImageSvgEditableConfig(
                    src=src,
                    out=out,
                    ocr_dir=ocr,
                    width=1920,
                    height=1080,
                    background="blank",
                    force=True,
                )
            )

            self.assertEqual(result.slide_count, 1)
            self.assertEqual(result.expected_text_blocks, 1)
            self.assertEqual(result.actual_text_boxes, 1)
            self.assertTrue((out / "slides_svg" / "slide_01.svg").exists())
            self.assertTrue((out / "editable_text_layer.pptx").exists())
            log = (out / "pipeline-execution-log.md").read_text(encoding="utf-8")
            self.assertIn("SVG output is a pixel-faithful wrapper", log)
            self.assertIn("Editable output is a separate PPTX", log)

    def test_blueprint_rebuild_creates_native_objects(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            blueprint = root / "deck.json"
            pptx = root / "blueprint.pptx"
            blueprint.write_text(
                json.dumps(
                    {
                        "canvas": {"width": 1920, "height": 1080},
                        "theme": {
                            "font": "Microsoft YaHei",
                            "accent": "#01e1d9",
                            "panel_fill": "#163f45",
                        },
                        "slides": [
                            {
                                "background": {"color": "#04191c"},
                                "elements": [
                                    {
                                        "type": "text",
                                        "x": 120,
                                        "y": 100,
                                        "w": 900,
                                        "h": 90,
                                        "text": "High fidelity title",
                                        "font_size": 28,
                                        "bold": True,
                                    },
                                    {
                                        "type": "component",
                                        "name": "panel",
                                        "x": 120,
                                        "y": 260,
                                        "w": 560,
                                        "h": 260,
                                        "title": "Native panel",
                                        "body": "Text, shape, and accent bar are editable.",
                                    },
                                    {
                                        "type": "component",
                                        "name": "chip",
                                        "x": 760,
                                        "y": 260,
                                        "w": 240,
                                        "h": 64,
                                        "text": "Editable chip",
                                    },
                                    {
                                        "type": "line",
                                        "x1": 760,
                                        "y1": 390,
                                        "x2": 1100,
                                        "y2": 390,
                                    },
                                    {
                                        "type": "component",
                                        "name": "footer",
                                        "page": "01",
                                        "label": "Blueprint",
                                    },
                                ],
                            }
                        ],
                    }
                ),
                encoding="utf-8",
            )

            result = rebuild_from_blueprint(
                BlueprintRebuildConfig(blueprint_path=blueprint, pptx_path=pptx)
            )

            self.assertTrue(pptx.exists())
            self.assertTrue(result.log_path.exists())
            self.assertEqual(result.slide_count, 1)
            self.assertGreaterEqual(result.text_count, 5)
            self.assertGreaterEqual(result.shape_count, 3)
            self.assertEqual(result.line_count, 1)

    def test_full_rebuild_pipeline_combines_exact_editable_and_blueprint_routes(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            src = root / "slides"
            ocr = root / "ocr"
            out = root / "out"
            src.mkdir()
            ocr.mkdir()
            (src / "slide_01.png").write_bytes(base64.b64decode(TINY_PNG))
            (ocr / "slide_01.json").write_text(
                '{"blocks":[{"text":"Hello","bbox":[0,0,100,40]}]}',
                encoding="utf-8",
            )
            blueprint = root / "deck.blueprint.json"
            blueprint.write_text(
                json.dumps(
                    {
                        "canvas": {"width": 1920, "height": 1080},
                        "slides": [
                            {
                                "background": {"color": "#ffffff"},
                                "elements": [
                                    {
                                        "type": "text",
                                        "x": 120,
                                        "y": 100,
                                        "w": 900,
                                        "h": 90,
                                        "text": "Native rebuild",
                                        "font_size": 28,
                                    }
                                ],
                            }
                        ],
                    }
                ),
                encoding="utf-8",
            )

            result = run_full_rebuild_pipeline(
                FullRebuildConfig(
                    src=src,
                    out=out,
                    ocr_dir=ocr,
                    blueprint_path=blueprint,
                    background="blank",
                    force=True,
                )
            )

            self.assertEqual(result.slide_count, 1)
            self.assertTrue((out / "exact_image_deck.pptx").exists())
            self.assertTrue((out / "editable_text_layer.pptx").exists())
            self.assertTrue((out / "high_fidelity_editable.pptx").exists())
            self.assertTrue((out / "full-rebuild-log.md").exists())
            self.assertIsNotNone(result.blueprint_pptx_path)
            self.assertEqual(result.expected_text_blocks, 1)
            self.assertEqual(result.actual_text_boxes, 1)
            log = (out / "full-rebuild-log.md").read_text(encoding="utf-8")
            self.assertIn("Exact proof route", log)
            self.assertIn("Editable text route", log)
            self.assertIn("High-fidelity blueprint route", log)

    def test_svg_native_rebuild_creates_editable_objects_from_structured_svg(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            src = root / "svg"
            src.mkdir()
            (src / "slide_01.svg").write_text(
                """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1920 1080">
  <rect x="100" y="120" width="520" height="240" fill="#163f45" stroke="#01e1d9" stroke-width="4"/>
  <text x="140" y="210" font-size="44" font-family="Microsoft YaHei" fill="#f4f8f7" font-weight="700">Native SVG title</text>
  <line x1="140" y1="260" x2="560" y2="260" stroke="#99efe9" stroke-width="3"/>
</svg>""",
                encoding="utf-8",
            )
            pptx = root / "native.pptx"

            result = rebuild_from_svg_native(
                SvgNativeRebuildConfig(src=src, pptx_path=pptx)
            )

            self.assertTrue(pptx.exists())
            self.assertTrue(result.log_path.exists())
            self.assertEqual(result.slide_count, 1)
            self.assertEqual(result.text_count, 1)
            self.assertGreaterEqual(result.shape_count, 1)
            self.assertEqual(result.line_count, 1)

            prs = Presentation(str(pptx))
            texts = [
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ]
            self.assertIn("Native SVG title", texts)

    def test_svg_native_rebuild_inserts_data_uri_images(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            src = root / "svg"
            src.mkdir()
            (src / "slide_01.svg").write_text(
                f"""<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1920 1080">
  <image x="10" y="20" width="100" height="80" href="data:image/png;base64,{TINY_PNG}"/>
</svg>""",
                encoding="utf-8",
            )
            pptx = root / "native.pptx"

            result = rebuild_from_svg_native(
                SvgNativeRebuildConfig(src=src, pptx_path=pptx)
            )

            self.assertEqual(result.picture_count, 1)
            prs = Presentation(str(pptx))
            pictures = [
                shape
                for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ]
            self.assertEqual(len(pictures), 1)

    def test_svg_native_rebuild_converts_polygon_and_simple_path(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            src = root / "svg"
            src.mkdir()
            (src / "slide_01.svg").write_text(
                """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1920 1080">
  <polygon points="120,120 240,120 200,220" fill="#01e1d9"/>
  <path d="M 320 120 L 460 120 L 460 220 Z" fill="#163f45" stroke="#99efe9"/>
</svg>""",
                encoding="utf-8",
            )
            pptx = root / "native.pptx"

            result = rebuild_from_svg_native(
                SvgNativeRebuildConfig(src=src, pptx_path=pptx)
            )

            self.assertEqual(result.shape_count, 2)
            self.assertEqual(result.skipped_count, 0)

    def test_cli_exposes_svg_native_rebuild_route(self) -> None:
        parser = build_parser()

        args = parser.parse_args(
            ["svg-native-rebuild", "svg", "--pptx", "native.pptx"]
        )

        self.assertEqual(args.command, "svg-native-rebuild")
        self.assertEqual(args.pattern, "slide_*.svg")

    def test_cli_exposes_ocr_filtering_arguments(self) -> None:
        parser = build_parser()

        editable_args = parser.parse_args(
            [
                "editable",
                "slides",
                "--ocr",
                "ocr",
                "--pptx",
                "editable.pptx",
                "--min-text-height",
                "30",
                "--min-text-area",
                "900",
                "--lock-file",
                "ocr-locks.json",
            ]
        )
        pipeline_args = parser.parse_args(
            [
                "image-svg-editable",
                "slides",
                "--out",
                "out",
                "--min-text-height",
                "30",
                "--min-text-area",
                "900",
                "--lock-file",
                "ocr-locks.json",
            ]
        )
        full_args = parser.parse_args(
            [
                "full-rebuild",
                "slides",
                "--out",
                "out",
                "--min-text-height",
                "30",
                "--min-text-area",
                "900",
                "--lock-file",
                "ocr-locks.json",
            ]
        )

        self.assertEqual(editable_args.min_text_height, 30)
        self.assertEqual(editable_args.min_text_area, 900)
        self.assertEqual(editable_args.lock_file, Path("ocr-locks.json"))
        self.assertEqual(pipeline_args.min_text_height, 30)
        self.assertEqual(pipeline_args.min_text_area, 900)
        self.assertEqual(pipeline_args.lock_file, Path("ocr-locks.json"))
        self.assertEqual(full_args.min_text_height, 30)
        self.assertEqual(full_args.min_text_area, 900)
        self.assertEqual(full_args.lock_file, Path("ocr-locks.json"))


if __name__ == "__main__":
    unittest.main()
