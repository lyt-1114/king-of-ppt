from __future__ import annotations

import base64
import html
import re
from dataclasses import dataclass
from pathlib import Path


MIME_BY_SUFFIX = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
}


@dataclass(frozen=True)
class ExportConfig:
    src: Path
    out: Path
    pattern: str = "slide_*.png"
    width: int = 1920
    height: int = 1080
    title: str = "Exact Image2PPT Export"
    pptx_path: Path | None = None
    create_pptx: bool = True
    force: bool = False


@dataclass(frozen=True)
class ExportResult:
    slide_count: int
    svg_dir: Path
    preview_path: Path
    run_log_path: Path
    pptx_path: Path | None


def export_exact_deck(config: ExportConfig) -> ExportResult:
    src = config.src.resolve()
    out = config.out.resolve()
    svg_dir = out / "slides_svg"

    images = collect_slide_images(src, config.pattern)
    if out.exists() and any(out.iterdir()) and not config.force:
        raise FileExistsError(
            f"Output folder is not empty: {out}. Pass --force to overwrite "
            "generated files."
        )

    svg_dir.mkdir(parents=True, exist_ok=True)
    generated_svgs: list[Path] = []
    for image_path in images:
        svg_path = svg_dir / f"{image_path.stem}.svg"
        svg_path.write_text(
            build_svg_wrapper(image_path, config.width, config.height),
            encoding="utf-8",
        )
        generated_svgs.append(svg_path)

    preview_path = out / "index.html"
    preview_path.write_text(
        build_preview_html(config.title, generated_svgs),
        encoding="utf-8",
    )

    pptx_path = None
    if config.create_pptx:
        pptx_path = (config.pptx_path or (out / "exact_image_deck.pptx")).resolve()
        create_image_pptx(images, pptx_path, config.width, config.height)

    run_log_path = out / "run-log.md"
    run_log_path.write_text(
        build_run_log(config, src, out, len(images), pptx_path),
        encoding="utf-8",
    )

    return ExportResult(
        slide_count=len(images),
        svg_dir=svg_dir,
        preview_path=preview_path,
        run_log_path=run_log_path,
        pptx_path=pptx_path,
    )


def collect_slide_images(src: Path, pattern: str) -> list[Path]:
    if not src.exists():
        raise FileNotFoundError(f"Source image folder not found: {src}")
    if not src.is_dir():
        raise NotADirectoryError(f"Source path is not a folder: {src}")

    images = [
        path
        for path in src.glob(pattern)
        if path.is_file() and path.suffix.lower() in MIME_BY_SUFFIX
    ]
    images.sort(key=natural_path_key)
    if not images:
        supported = ", ".join(sorted(MIME_BY_SUFFIX))
        raise FileNotFoundError(
            f"No supported slide images found under {src} with pattern "
            f"{pattern!r}. Supported extensions: {supported}"
        )
    return images


def natural_path_key(path: Path) -> list[object]:
    parts = re.split(r"(\d+)", path.name.lower())
    return [int(part) if part.isdigit() else part for part in parts]


def build_svg_wrapper(image_path: Path, width: int, height: int) -> str:
    data_uri = image_to_data_uri(image_path)
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="no"?>
<svg xmlns="http://www.w3.org/2000/svg"
     xmlns:xlink="http://www.w3.org/1999/xlink"
     width="{width}"
     height="{height}"
     viewBox="0 0 {width} {height}"
     version="1.1">
  <title>Exact slide reproduction</title>
  <desc>Pixel-identical SVG wrapper generated from a full-slide image render.</desc>
  <image x="0" y="0" width="{width}" height="{height}" preserveAspectRatio="none" xlink:href="{data_uri}" href="{data_uri}" />
</svg>
"""


def image_to_data_uri(path: Path) -> str:
    suffix = path.suffix.lower()
    mime = MIME_BY_SUFFIX.get(suffix)
    if mime is None:
        raise ValueError(f"Unsupported image type: {path}")
    encoded = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:{mime};base64,{encoded}"


def build_preview_html(title: str, slides: list[Path]) -> str:
    safe_title = html.escape(title)
    body = "\n".join(
        f"""<section class="slide">
  <div class="label">{html.escape(slide.name)}</div>
  <img src="slides_svg/{html.escape(slide.name)}" alt="{html.escape(slide.name)}" />
</section>"""
        for slide in slides
    )
    return f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{safe_title}</title>
  <style>
    :root {{
      color-scheme: dark;
      --bg: #081317;
      --panel: #10272c;
      --text: #e7f1ef;
      --muted: #98abad;
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      background: var(--bg);
      color: var(--text);
      font-family: "Segoe UI", Arial, sans-serif;
      padding: 24px;
    }}
    .wrap {{
      max-width: 1400px;
      margin: 0 auto;
      display: grid;
      gap: 24px;
    }}
    .note {{
      background: var(--panel);
      border: 1px solid rgba(160, 230, 220, 0.15);
      border-radius: 14px;
      padding: 16px 18px;
      color: var(--muted);
      line-height: 1.6;
    }}
    .slide {{ display: grid; gap: 10px; }}
    .label {{ color: var(--muted); font-size: 14px; }}
    img {{
      width: 100%;
      height: auto;
      display: block;
      border-radius: 12px;
      border: 1px solid rgba(160, 230, 220, 0.15);
      background: #000;
    }}
  </style>
</head>
<body>
  <div class="wrap">
    <div class="note">
      This preview uses SVG files that embed the approved slide images as
      full-canvas bitmap layers. It is built for pixel-level reproduction, not
      editable object recovery.
    </div>
    {body}
  </div>
</body>
</html>
"""


def create_image_pptx(
    images: list[Path],
    pptx_path: Path,
    width: int,
    height: int,
) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is required to create PPTX output. Install this "
            "package with pip install -e ."
        ) from exc

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(width / 144)
    presentation.slide_height = Inches(height / 144)
    blank_layout = presentation.slide_layouts[6]

    for image_path in images:
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

    presentation.save(pptx_path)


def build_run_log(
    config: ExportConfig,
    src: Path,
    out: Path,
    slide_count: int,
    pptx_path: Path | None,
) -> str:
    lines = [
        "# Exact Image2PPT Export",
        "",
        f"- Source slide renders: `{src}`",
        f"- Output folder: `{out}`",
        f"- Slide count: `{slide_count}`",
        f"- Canvas: `{config.width} x {config.height}`",
        f"- Pattern: `{config.pattern}`",
        "- SVG strategy: each SVG embeds the source slide image as a full-canvas data URI.",
        "- PPTX strategy: each slide uses the source image as a full-canvas picture layer.",
        "- Fidelity target: pixel-level visual reproduction of the approved image render.",
        "- Editability limit: text, shapes, charts, and diagrams are not recovered as native PPT objects.",
    ]
    if pptx_path is not None:
        lines.append(f"- PPTX output: `{pptx_path}`")
    return "\n".join(lines) + "\n"
