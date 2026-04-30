from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation

from .editable import (
    EditableConfig,
    OcrConfig,
    create_editable_text_pptx_with_stats,
    extract_ocr_json,
    load_slide_blocks,
)
from .exporter import ExportConfig, collect_slide_images, export_exact_deck


@dataclass(frozen=True)
class ImageSvgEditableConfig:
    src: Path
    out: Path
    editable_pptx_path: Path | None = None
    ocr_dir: Path | None = None
    pattern: str = "slide_*.png"
    width: int = 1920
    height: int = 1080
    lang: str = "ch"
    background: str = "blank"
    default_font: str = "Microsoft YaHei"
    default_color: str = "#111827"
    min_text_height: float | None = None
    min_text_area: float | None = None
    lock_file: Path | None = None
    force: bool = False
    allow_empty_text: bool = False


@dataclass(frozen=True)
class ImageSvgEditableResult:
    slide_count: int
    svg_dir: Path
    preview_path: Path
    exact_pptx_path: Path | None
    ocr_dir: Path
    editable_pptx_path: Path
    pipeline_log_path: Path
    expected_text_blocks: int
    actual_text_boxes: int
    skipped_by_height: int = 0
    skipped_by_area: int = 0
    skipped_by_lock: int = 0


def run_image_svg_editable_pipeline(
    config: ImageSvgEditableConfig,
) -> ImageSvgEditableResult:
    out = config.out.resolve()
    editable_pptx_path = (
        config.editable_pptx_path or out / "editable_text_layer.pptx"
    ).resolve()

    exact_result = export_exact_deck(
        ExportConfig(
            src=config.src,
            out=out,
            pattern=config.pattern,
            width=config.width,
            height=config.height,
            title="Image SVG Editable Pipeline",
            pptx_path=out / "exact_image_deck.pptx",
            create_pptx=True,
            force=config.force,
        )
    )

    ocr_dir = (config.ocr_dir or out / "ocr_json").resolve()
    if config.ocr_dir is None:
        extract_ocr_json(
            OcrConfig(
                src=config.src,
                out=ocr_dir,
                pattern=config.pattern,
                lang=config.lang,
            )
        )

    expected_blocks = count_ocr_blocks_for_slides(
        config.src.resolve(),
        ocr_dir,
        config.pattern,
    )
    if expected_blocks == 0 and not config.allow_empty_text:
        raise RuntimeError(
            "OCR JSON contains zero text blocks. Refusing to create a falsely "
            "successful editable PPTX. Provide corrected OCR JSON or pass "
            "--allow-empty-text for image-only decks."
        )

    editable_result = create_editable_text_pptx_with_stats(
        EditableConfig(
            src=config.src,
            ocr_dir=ocr_dir,
            pptx_path=editable_pptx_path,
            pattern=config.pattern,
            width=config.width,
            height=config.height,
            background=config.background,
            default_font=config.default_font,
            default_color=config.default_color,
            min_text_height=config.min_text_height,
            min_text_area=config.min_text_area,
            lock_file=config.lock_file,
        )
    )

    slide_count, actual_text_boxes = inspect_editable_pptx(editable_pptx_path)
    if slide_count != exact_result.slide_count:
        raise RuntimeError(
            f"Editable PPTX slide count mismatch: expected "
            f"{exact_result.slide_count}, got {slide_count}."
        )
    if actual_text_boxes == 0 and not config.allow_empty_text:
        raise RuntimeError(
            "Editable PPTX contains zero editable text boxes after generation."
        )

    pipeline_log_path = out / "pipeline-execution-log.md"
    pipeline_log_path.write_text(
        build_pipeline_execution_log(
            config=config,
            slide_count=exact_result.slide_count,
            svg_dir=exact_result.svg_dir,
            preview_path=exact_result.preview_path,
            exact_pptx_path=exact_result.pptx_path,
            ocr_dir=ocr_dir,
            editable_pptx_path=editable_pptx_path,
            expected_text_blocks=expected_blocks,
            actual_text_boxes=actual_text_boxes,
            skipped_by_height=editable_result.skipped_by_height,
            skipped_by_area=editable_result.skipped_by_area,
            skipped_by_lock=editable_result.skipped_by_lock,
        ),
        encoding="utf-8",
    )

    return ImageSvgEditableResult(
        slide_count=exact_result.slide_count,
        svg_dir=exact_result.svg_dir,
        preview_path=exact_result.preview_path,
        exact_pptx_path=exact_result.pptx_path,
        ocr_dir=ocr_dir,
        editable_pptx_path=editable_pptx_path,
        pipeline_log_path=pipeline_log_path,
        expected_text_blocks=expected_blocks,
        actual_text_boxes=actual_text_boxes,
        skipped_by_height=editable_result.skipped_by_height,
        skipped_by_area=editable_result.skipped_by_area,
        skipped_by_lock=editable_result.skipped_by_lock,
    )


def count_ocr_blocks_for_slides(src: Path, ocr_dir: Path, pattern: str) -> int:
    total = 0
    missing: list[str] = []
    for image_path in collect_slide_images(src, pattern):
        json_path = ocr_dir / f"{image_path.stem}.json"
        if not json_path.exists():
            missing.append(json_path.name)
            continue
        total += len(load_slide_blocks(ocr_dir, image_path.stem))
    if missing:
        raise FileNotFoundError(
            "Missing OCR JSON files for slides: " + ", ".join(missing)
        )
    return total


def inspect_editable_pptx(path: Path) -> tuple[int, int]:
    prs = Presentation(str(path))
    text_boxes = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_boxes += 1
    return len(prs.slides), text_boxes


def build_pipeline_execution_log(
    *,
    config: ImageSvgEditableConfig,
    slide_count: int,
    svg_dir: Path,
    preview_path: Path,
    exact_pptx_path: Path | None,
    ocr_dir: Path,
    editable_pptx_path: Path,
    expected_text_blocks: int,
    actual_text_boxes: int,
    skipped_by_height: int,
    skipped_by_area: int,
    skipped_by_lock: int,
) -> str:
    return "\n".join(
        [
            "# Image SVG Editable Pipeline Execution Log",
            "",
            "## Route",
            "",
            "- Source images are treated as the approved visual reference.",
            "- SVG output is a pixel-faithful wrapper: each SVG embeds one full-slide bitmap image.",
            "- SVG-in-PPT does not mean PPT editability when the SVG itself is one embedded image.",
            "- Editable output is a separate PPTX rebuilt with native PowerPoint text boxes.",
            "- The source image deck is never overwritten by this pipeline.",
            "",
            "## Tradeoff",
            "",
            "- Pixel-identical reproduction is provided by the exact image/SVG/PPTX route.",
            "- Editable PPTX is provided by reconstructing native objects from OCR JSON.",
            "- A deck cannot be both 100 percent pixel-identical to a flattened bitmap and fully editable as native PPT objects at the same time.",
            "- For higher visual fidelity, enrich the editable reconstruction with native shapes, lines, diagrams, and replaceable image/SVG decorations.",
            "- Background mode: `" + config.background + "`.",
            "- Warning: `keep` preserves source slide images under OCR text. Use it only when the source images are text-free or when creating a visual debug layer; otherwise duplicate text is expected.",
            "",
            "## Inputs",
            "",
            f"- Source slide images: `{config.src.resolve()}`",
            f"- Pattern: `{config.pattern}`",
            f"- Canvas: `{config.width} x {config.height}`",
            f"- OCR JSON: `{ocr_dir}`",
            f"- Minimum editable text height: `{config.min_text_height if config.min_text_height is not None else 'disabled'}`",
            f"- Minimum editable text area: `{config.min_text_area if config.min_text_area is not None else 'disabled'}`",
            f"- OCR lock file: `{config.lock_file.resolve() if config.lock_file else 'not provided'}`",
            "",
            "## Outputs",
            "",
            f"- Slide count: `{slide_count}`",
            f"- SVG wrappers: `{svg_dir}`",
            f"- SVG preview: `{preview_path}`",
            f"- Exact image PPTX: `{exact_pptx_path}`",
            f"- Editable PPTX: `{editable_pptx_path}`",
            "",
            "## Verification",
            "",
            f"- OCR text blocks found: `{expected_text_blocks}`",
            f"- Editable PPT text boxes found: `{actual_text_boxes}`",
            f"- OCR blocks skipped by height: `{skipped_by_height}`",
            f"- OCR blocks skipped by area: `{skipped_by_area}`",
            f"- OCR blocks skipped by lock regions: `{skipped_by_lock}`",
            "- Status: editable text layer exists and passed slide-count/text-box checks.",
            "",
        ]
    )
