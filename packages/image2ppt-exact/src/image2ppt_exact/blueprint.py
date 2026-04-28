from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


@dataclass(frozen=True)
class BlueprintRebuildConfig:
    blueprint_path: Path
    pptx_path: Path
    assets_root: Path | None = None


@dataclass(frozen=True)
class BlueprintRebuildResult:
    pptx_path: Path
    log_path: Path
    slide_count: int
    text_count: int
    shape_count: int
    picture_count: int
    line_count: int


def rebuild_from_blueprint(config: BlueprintRebuildConfig) -> BlueprintRebuildResult:
    blueprint_path = config.blueprint_path.resolve()
    data = json.loads(blueprint_path.read_text(encoding="utf-8-sig"))
    assets_root = (config.assets_root or blueprint_path.parent).resolve()

    canvas = data.get("canvas", {})
    source_w = int(canvas.get("width", 1920))
    source_h = int(canvas.get("height", 1080))
    slide_w_in = float(canvas.get("slide_width_in", source_w / 144))
    slide_h_in = float(canvas.get("slide_height_in", source_h / 144))

    prs = Presentation()
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)
    blank_layout = prs.slide_layouts[6]

    stats = {"text": 0, "shape": 0, "picture": 0, "line": 0}
    theme = data.get("theme", {})

    for slide_data in data.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        if "background" in slide_data:
            add_background(slide, prs, slide_data["background"], source_w, source_h, assets_root)
        for element in slide_data.get("elements", []):
            add_element(slide, prs, element, source_w, source_h, assets_root, theme, stats)

    config.pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(config.pptx_path)

    log_path = config.pptx_path.with_suffix(".blueprint-log.md")
    log_path.write_text(
        build_blueprint_log(
            blueprint_path=blueprint_path,
            assets_root=assets_root,
            pptx_path=config.pptx_path,
            slide_count=len(prs.slides),
            stats=stats,
        ),
        encoding="utf-8",
    )

    return BlueprintRebuildResult(
        pptx_path=config.pptx_path,
        log_path=log_path,
        slide_count=len(prs.slides),
        text_count=stats["text"],
        shape_count=stats["shape"],
        picture_count=stats["picture"],
        line_count=stats["line"],
    )


def add_background(
    slide: Any,
    prs: Presentation,
    background: Any,
    source_w: int,
    source_h: int,
    assets_root: Path,
) -> None:
    if isinstance(background, str):
        add_full_slide_picture(slide, prs, resolve_asset(background, assets_root))
        return
    if not isinstance(background, dict):
        return
    if "image" in background:
        add_full_slide_picture(slide, prs, resolve_asset(background["image"], assets_root))
    if "color" in background:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            prs.slide_width,
            prs.slide_height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = parse_color(background["color"])
        shape.line.fill.background()


def add_full_slide_picture(slide: Any, prs: Presentation, path: Path) -> None:
    slide.shapes.add_picture(
        str(path),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )


def add_element(
    slide: Any,
    prs: Presentation,
    element: dict[str, Any],
    source_w: int,
    source_h: int,
    assets_root: Path,
    theme: dict[str, Any],
    stats: dict[str, int],
) -> None:
    element_type = element.get("type")
    if element_type == "picture":
        add_picture(slide, prs, element, source_w, source_h, assets_root)
        stats["picture"] += 1
        return
    if element_type == "text":
        add_textbox(slide, prs, element, source_w, source_h, theme)
        stats["text"] += 1
        return
    if element_type == "shape":
        shape = add_shape(slide, prs, element, source_w, source_h)
        stats["shape"] += 1
        if element.get("text"):
            set_text(shape, element, theme)
            stats["text"] += 1
        return
    if element_type == "line":
        add_line(slide, prs, element, source_w, source_h)
        stats["line"] += 1
        return
    if element_type == "component":
        add_component(slide, prs, element, source_w, source_h, theme, stats)
        return
    raise ValueError(f"Unsupported blueprint element type: {element_type!r}")


def add_picture(
    slide: Any,
    prs: Presentation,
    element: dict[str, Any],
    source_w: int,
    source_h: int,
    assets_root: Path,
) -> None:
    x, y, w, h = box(element, source_w, source_h, prs)
    slide.shapes.add_picture(
        str(resolve_asset(element["path"], assets_root)),
        x,
        y,
        width=w,
        height=h,
    )


def add_textbox(
    slide: Any,
    prs: Presentation,
    element: dict[str, Any],
    source_w: int,
    source_h: int,
    theme: dict[str, Any],
) -> Any:
    x, y, w, h = box(element, source_w, source_h, prs)
    shape = slide.shapes.add_textbox(x, y, w, h)
    set_text(shape, element, theme)
    return shape


def add_shape(
    slide: Any,
    prs: Presentation,
    element: dict[str, Any],
    source_w: int,
    source_h: int,
) -> Any:
    x, y, w, h = box(element, source_w, source_h, prs)
    shape = slide.shapes.add_shape(shape_type(element.get("shape", "rect")), x, y, w, h)
    if "fill" in element:
        shape.fill.solid()
        shape.fill.fore_color.rgb = parse_color(element["fill"])
    else:
        shape.fill.background()
    if "line" in element:
        line = element["line"] or {}
        if isinstance(line, str):
            shape.line.color.rgb = parse_color(line)
        else:
            shape.line.color.rgb = parse_color(line.get("color", "#000000"))
            if "width" in line:
                shape.line.width = Pt(float(line["width"]))
    else:
        shape.line.fill.background()
    return shape


def add_line(
    slide: Any,
    prs: Presentation,
    element: dict[str, Any],
    source_w: int,
    source_h: int,
) -> None:
    x1 = px_to_emu(float(element["x1"]), source_w, prs.slide_width)
    y1 = px_to_emu(float(element["y1"]), source_h, prs.slide_height)
    x2 = px_to_emu(float(element["x2"]), source_w, prs.slide_width)
    y2 = px_to_emu(float(element["y2"]), source_h, prs.slide_height)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = parse_color(element.get("color", "#99efe9"))
    line.line.width = Pt(float(element.get("width", 1.2)))


def add_component(
    slide: Any,
    prs: Presentation,
    element: dict[str, Any],
    source_w: int,
    source_h: int,
    theme: dict[str, Any],
    stats: dict[str, int],
) -> None:
    name = element.get("name")
    if name == "panel":
        add_panel(slide, prs, element, source_w, source_h, theme, stats)
        return
    if name == "chip":
        add_chip(slide, prs, element, source_w, source_h, theme, stats)
        return
    if name == "footer":
        add_footer(slide, prs, element, source_w, source_h, theme, stats)
        return
    raise ValueError(f"Unsupported blueprint component: {name!r}")


def add_panel(
    slide: Any,
    prs: Presentation,
    element: dict[str, Any],
    source_w: int,
    source_h: int,
    theme: dict[str, Any],
    stats: dict[str, int],
) -> None:
    x, y, w, h = element["x"], element["y"], element["w"], element["h"]
    fill = element.get("fill", theme.get("panel_fill", "#163f45"))
    accent = element.get("accent", theme.get("accent", "#01e1d9"))
    add_element(
        slide,
        prs,
        {"type": "shape", "shape": "roundRect", "x": x, "y": y, "w": w, "h": h, "fill": fill, "line": {"color": accent, "width": 1}},
        source_w,
        source_h,
        Path("."),
        theme,
        stats,
    )
    add_element(
        slide,
        prs,
        {"type": "shape", "shape": "rect", "x": x + 16, "y": y + 18, "w": 8, "h": max(12, h - 36), "fill": accent},
        source_w,
        source_h,
        Path("."),
        theme,
        stats,
    )
    add_element(
        slide,
        prs,
        {"type": "text", "x": x + 42, "y": y + 22, "w": w - 58, "h": 42, "text": element.get("title", ""), "font_size": element.get("title_size", 14), "bold": True, "color": element.get("title_color", "#f4f8f7")},
        source_w,
        source_h,
        Path("."),
        theme,
        stats,
    )
    add_element(
        slide,
        prs,
        {"type": "text", "x": x + 42, "y": y + 72, "w": w - 58, "h": max(24, h - 86), "text": element.get("body", ""), "font_size": element.get("body_size", 12), "color": element.get("body_color", "#cad7d8")},
        source_w,
        source_h,
        Path("."),
        theme,
        stats,
    )


def add_chip(
    slide: Any,
    prs: Presentation,
    element: dict[str, Any],
    source_w: int,
    source_h: int,
    theme: dict[str, Any],
    stats: dict[str, int],
) -> None:
    x, y, w, h = element["x"], element["y"], element["w"], element["h"]
    accent = element.get("accent", theme.get("accent", "#01e1d9"))
    add_element(
        slide,
        prs,
        {"type": "shape", "shape": "roundRect", "x": x, "y": y, "w": w, "h": h, "fill": element.get("fill", "#11383e"), "line": {"color": accent, "width": 1.2}},
        source_w,
        source_h,
        Path("."),
        theme,
        stats,
    )
    add_element(
        slide,
        prs,
        {"type": "text", "x": x + 4, "y": y + 4, "w": w - 8, "h": h - 8, "text": element.get("text", ""), "font_size": element.get("font_size", 12), "bold": True, "align": "center", "color": element.get("color", "#f4f8f7")},
        source_w,
        source_h,
        Path("."),
        theme,
        stats,
    )


def add_footer(
    slide: Any,
    prs: Presentation,
    element: dict[str, Any],
    source_w: int,
    source_h: int,
    theme: dict[str, Any],
    stats: dict[str, int],
) -> None:
    color = element.get("color", theme.get("muted", "#8ea6a8"))
    accent = element.get("accent", theme.get("accent", "#01e1d9"))
    y = element.get("y", 1022)
    add_element(
        slide,
        prs,
        {"type": "text", "x": 72, "y": y, "w": 80, "h": 30, "text": str(element.get("page", "")), "font_size": 10.5, "bold": True, "color": accent},
        source_w,
        source_h,
        Path("."),
        theme,
        stats,
    )
    add_element(
        slide,
        prs,
        {"type": "text", "x": 126, "y": y, "w": 410, "h": 32, "text": element.get("label", ""), "font_size": 10.2, "color": color},
        source_w,
        source_h,
        Path("."),
        theme,
        stats,
    )
    if element.get("source"):
        add_element(
            slide,
            prs,
            {"type": "text", "x": 880, "y": y, "w": 930, "h": 32, "text": element["source"], "font_size": 9.2, "align": "right", "color": color},
            source_w,
            source_h,
            Path("."),
            theme,
            stats,
        )


def set_text(shape: Any, element: dict[str, Any], theme: dict[str, Any]) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = vertical_anchor(element.get("valign", "top"))
    margin = float(element.get("margin", 3))
    frame.margin_left = Pt(margin)
    frame.margin_right = Pt(margin)
    frame.margin_top = Pt(margin)
    frame.margin_bottom = Pt(margin)

    lines = str(element.get("text", "")).split("\n")
    for idx, text in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.alignment = alignment(element.get("align", "left"))
        paragraph.line_spacing = float(element.get("line_spacing", 1.0))
        run = paragraph.add_run()
        run.text = text
        run.font.name = element.get("font", theme.get("font", "Microsoft YaHei"))
        run.font.size = Pt(float(element.get("font_size", theme.get("font_size", 12))))
        run.font.bold = bool(element.get("bold", False))
        run.font.italic = bool(element.get("italic", False))
        run.font.color.rgb = parse_color(element.get("color", theme.get("text", "#f4f8f7")))


def box(element: dict[str, Any], source_w: int, source_h: int, prs: Presentation) -> tuple[int, int, int, int]:
    return (
        px_to_emu(float(element["x"]), source_w, prs.slide_width),
        px_to_emu(float(element["y"]), source_h, prs.slide_height),
        px_to_emu(float(element["w"]), source_w, prs.slide_width),
        px_to_emu(float(element["h"]), source_h, prs.slide_height),
    )


def px_to_emu(value: float, source_extent: int, slide_extent: int) -> int:
    return int(round((value / source_extent) * slide_extent))


def resolve_asset(path_value: str, assets_root: Path) -> Path:
    path = Path(path_value)
    if path.is_absolute():
        return path
    return (assets_root / path).resolve()


def parse_color(value: str | RGBColor) -> RGBColor:
    if isinstance(value, RGBColor):
        return value
    raw = str(value).strip().lstrip("#")
    if len(raw) != 6:
        raise ValueError(f"Expected #RRGGBB color, got {value!r}")
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def shape_type(value: str) -> Any:
    normalized = value.lower()
    if normalized in {"roundrect", "rounded_rect", "rounded_rectangle"}:
        return MSO_SHAPE.ROUNDED_RECTANGLE
    if normalized == "oval":
        return MSO_SHAPE.OVAL
    return MSO_SHAPE.RECTANGLE


def alignment(value: str) -> PP_ALIGN:
    normalized = str(value).lower()
    if normalized == "center":
        return PP_ALIGN.CENTER
    if normalized == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def vertical_anchor(value: str) -> MSO_VERTICAL_ANCHOR:
    normalized = str(value).lower()
    if normalized == "middle":
        return MSO_VERTICAL_ANCHOR.MIDDLE
    if normalized == "bottom":
        return MSO_VERTICAL_ANCHOR.BOTTOM
    return MSO_VERTICAL_ANCHOR.TOP


def build_blueprint_log(
    *,
    blueprint_path: Path,
    assets_root: Path,
    pptx_path: Path,
    slide_count: int,
    stats: dict[str, int],
) -> str:
    return "\n".join(
        [
            "# Blueprint Rebuild Log",
            "",
            "- Route: blueprint JSON -> native PowerPoint objects.",
            "- Purpose: high-fidelity editable reconstruction using text boxes, shapes, lines, pictures, and reusable components.",
            "- This is the route that matches the successful high-fidelity editable deck pattern better than OCR-only reconstruction.",
            "",
            f"- Blueprint: `{blueprint_path}`",
            f"- Assets root: `{assets_root}`",
            f"- Output PPTX: `{pptx_path.resolve()}`",
            f"- Slides: `{slide_count}`",
            f"- Native text objects: `{stats['text']}`",
            f"- Native shape objects: `{stats['shape']}`",
            f"- Picture objects: `{stats['picture']}`",
            f"- Line objects: `{stats['line']}`",
            "",
        ]
    )
