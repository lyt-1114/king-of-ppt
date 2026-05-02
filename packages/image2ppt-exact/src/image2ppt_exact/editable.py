from __future__ import annotations

import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from difflib import SequenceMatcher

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from .exporter import collect_slide_images


@dataclass(frozen=True)
class EditableConfig:
    src: Path
    ocr_dir: Path
    pptx_path: Path
    pattern: str = "slide_*.png"
    width: int = 1920
    height: int = 1080
    background: str = "blank"
    default_font: str = "Microsoft YaHei"
    default_color: str = "#111827"
    min_text_height: float | None = None
    min_text_area: float | None = None
    lock_file: Path | None = None
    spec_file: Path | None = None


@dataclass(frozen=True)
class EditableFilterStats:
    total_blocks: int = 0
    editable_blocks: int = 0
    skipped_by_height: int = 0
    skipped_by_area: int = 0
    skipped_by_lock: int = 0


@dataclass(frozen=True)
class EditableBuildResult:
    pptx_path: Path
    total_blocks: int
    editable_blocks: int
    skipped_by_height: int
    skipped_by_area: int
    skipped_by_lock: int


@dataclass(frozen=True)
class OcrConfig:
    src: Path
    out: Path
    pattern: str = "slide_*.png"
    lang: str = "ch"


def create_editable_text_pptx(config: EditableConfig) -> Path:
    return create_editable_text_pptx_with_stats(config).pptx_path


def create_editable_text_pptx_with_stats(config: EditableConfig) -> EditableBuildResult:
    from pptx import Presentation
    from pptx.util import Inches

    images = collect_slide_images(config.src.resolve(), config.pattern)
    ocr_dir = config.ocr_dir.resolve()
    if not ocr_dir.exists():
        raise FileNotFoundError(f"OCR JSON folder not found: {ocr_dir}")

    prs = Presentation()
    prs.slide_width = Inches(config.width / 144)
    prs.slide_height = Inches(config.height / 144)
    blank_layout = prs.slide_layouts[6]
    default_color = parse_hex_color(config.default_color)
    locks = load_ocr_locks(config.lock_file)
    total_stats = EditableFilterStats()

    for image_path in images:
        slide = prs.slides.add_slide(blank_layout)
        blocks = load_slide_blocks(ocr_dir, image_path.stem)
        blocks = apply_spec_corrections(
            blocks,
            load_spec_lines(config.spec_file) if config.spec_file else [],
        )
        editable_blocks, stats = collect_editable_blocks(
            blocks,
            slide_stem=image_path.stem,
            width=config.width,
            height=config.height,
            min_text_height=config.min_text_height,
            min_text_area=config.min_text_area,
            locks=locks,
        )
        total_stats = combine_filter_stats(total_stats, stats)
        if config.background == "keep":
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
        elif config.background == "redact":
            redacted_path = create_redacted_background(
                image_path=image_path,
                blocks=editable_blocks,
                output_dir=config.pptx_path.with_suffix("").parent
                / f"{config.pptx_path.stem}_assets",
                width=config.width,
                height=config.height,
            )
            slide.shapes.add_picture(
                str(redacted_path),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
        elif config.background != "blank":
            raise ValueError("background must be 'blank', 'keep', or 'redact'")

        for block in editable_blocks:
            x, y, w, h = normalize_bbox(block["bbox"], config.width, config.height)
            shape = slide.shapes.add_textbox(
                px_to_emu(x, config.width, prs.slide_width),
                px_to_emu(y, config.height, prs.slide_height),
                px_to_emu(w, config.width, prs.slide_width),
                px_to_emu(h, config.height, prs.slide_height),
            )
            text_frame = shape.text_frame
            text_frame.clear()
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            text_frame.word_wrap = True

            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = parse_alignment(block.get("align"))
            run = paragraph.add_run()
            run.text = str(block.get("text", ""))
            run.font.name = str(block.get("font", config.default_font))
            run.font.size = Pt(
                float(block.get("font_size") or estimate_font_size(h, config.height))
            )
            run.font.bold = bool(block.get("bold", False))
            run.font.italic = bool(block.get("italic", False))
            run.font.color.rgb = (
                parse_hex_color(block.get("color", config.default_color))
                or default_color
            )

    config.pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(config.pptx_path)
    return EditableBuildResult(
        pptx_path=config.pptx_path,
        total_blocks=total_stats.total_blocks,
        editable_blocks=total_stats.editable_blocks,
        skipped_by_height=total_stats.skipped_by_height,
        skipped_by_area=total_stats.skipped_by_area,
        skipped_by_lock=total_stats.skipped_by_lock,
    )


def collect_editable_blocks(
    blocks: list[dict[str, Any]],
    *,
    slide_stem: str,
    width: int,
    height: int,
    min_text_height: float | None = None,
    min_text_area: float | None = None,
    locks: dict[str, list[dict[str, float]]] | None = None,
) -> tuple[list[dict[str, Any]], EditableFilterStats]:
    editable: list[dict[str, Any]] = []
    skipped_by_height = 0
    skipped_by_area = 0
    skipped_by_lock = 0
    slide_locks = (locks or {}).get("*", []) + (locks or {}).get(slide_stem, [])

    for block in blocks:
        x, y, w, h = normalize_bbox(block["bbox"], width, height)
        if min_text_height is not None and h < min_text_height:
            skipped_by_height += 1
            continue
        if min_text_area is not None and w * h < min_text_area:
            skipped_by_area += 1
            continue
        if any(boxes_intersect((x, y, w, h), lock) for lock in slide_locks):
            skipped_by_lock += 1
            continue
        editable.append(block)

    return editable, EditableFilterStats(
        total_blocks=len(blocks),
        editable_blocks=len(editable),
        skipped_by_height=skipped_by_height,
        skipped_by_area=skipped_by_area,
        skipped_by_lock=skipped_by_lock,
    )


def combine_filter_stats(
    left: EditableFilterStats,
    right: EditableFilterStats,
) -> EditableFilterStats:
    return EditableFilterStats(
        total_blocks=left.total_blocks + right.total_blocks,
        editable_blocks=left.editable_blocks + right.editable_blocks,
        skipped_by_height=left.skipped_by_height + right.skipped_by_height,
        skipped_by_area=left.skipped_by_area + right.skipped_by_area,
        skipped_by_lock=left.skipped_by_lock + right.skipped_by_lock,
    )


def boxes_intersect(
    bbox: tuple[float, float, float, float],
    lock: dict[str, float],
) -> bool:
    x, y, w, h = bbox
    ax1, ay1, ax2, ay2 = x, y, x + w, y + h
    bx1 = lock["x"]
    by1 = lock["y"]
    bx2 = bx1 + lock["w"]
    by2 = by1 + lock["h"]
    return min(ax2, bx2) > max(ax1, bx1) and min(ay2, by2) > max(ay1, by1)


def load_ocr_locks(lock_file: Path | None) -> dict[str, list[dict[str, float]]]:
    if lock_file is None:
        return {}
    data = json.loads(lock_file.read_text(encoding="utf-8-sig"))
    locks: dict[str, list[dict[str, float]]] = {}
    if isinstance(data, list):
        locks["*"] = [normalize_lock_region(region) for region in data]
        return locks
    if not isinstance(data, dict):
        raise ValueError("OCR lock file must contain a list or object.")

    global_regions = data.get("regions", [])
    if global_regions:
        if not isinstance(global_regions, list):
            raise ValueError("OCR lock file 'regions' must be a list.")
        locks["*"] = [normalize_lock_region(region) for region in global_regions]

    slide_regions = data.get("slides", {})
    if slide_regions:
        if not isinstance(slide_regions, dict):
            raise ValueError("OCR lock file 'slides' must be an object.")
        for slide_name, regions in slide_regions.items():
            if not isinstance(regions, list):
                raise ValueError(f"OCR lock regions for {slide_name!r} must be a list.")
            locks[str(slide_name)] = [normalize_lock_region(region) for region in regions]
    return locks


def normalize_lock_region(region: Any) -> dict[str, float]:
    if isinstance(region, dict):
        x = float(region.get("x", region.get("left", 0)))
        y = float(region.get("y", region.get("top", 0)))
        w = float(region.get("w", region.get("width", 0)))
        h = float(region.get("h", region.get("height", 0)))
    elif isinstance(region, list) and len(region) == 4:
        x, y, w, h = [float(value) for value in region]
    else:
        raise ValueError(f"Unsupported OCR lock region: {region!r}")
    if w <= 0 or h <= 0:
        raise ValueError(f"OCR lock region must have positive width and height: {region!r}")
    return {"x": x, "y": y, "w": w, "h": h}


def create_redacted_background(
    *,
    image_path: Path,
    blocks: list[dict[str, Any]],
    output_dir: Path,
    width: int,
    height: int,
    padding: int = 2,
) -> Path:
    try:
        from PIL import Image, ImageDraw
    except ImportError as exc:
        raise RuntimeError(
            "Pillow is required for background='redact'. Install Pillow or use "
            "background='blank'."
        ) from exc

    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"{image_path.stem}_redacted.png"
    with Image.open(image_path) as source:
        image = source.convert("RGB")
    draw = ImageDraw.Draw(image)

    scale_x = image.width / width
    scale_y = image.height / height
    for block in blocks:
        x, y, w, h = normalize_bbox(block["bbox"], width, height)
        left = max(0, int(round((x - padding) * scale_x)))
        top = max(0, int(round((y - padding) * scale_y)))
        right = min(image.width, int(round((x + w + padding) * scale_x)))
        bottom = min(image.height, int(round((y + h + padding) * scale_y)))
        fill = sample_background_color(image, left, top, right, bottom)
        draw.rectangle((left, top, right, bottom), fill=fill)

    image.save(output_path)
    return output_path


def sample_background_color(
    image: Any,
    left: int,
    top: int,
    right: int,
    bottom: int,
) -> tuple[int, int, int]:
    samples: list[tuple[int, int, int]] = []
    candidates = [
        (left - 2, top - 2),
        (right + 2, top - 2),
        (left - 2, bottom + 2),
        (right + 2, bottom + 2),
        ((left + right) // 2, top - 2),
        ((left + right) // 2, bottom + 2),
        (left - 2, (top + bottom) // 2),
        (right + 2, (top + bottom) // 2),
    ]
    for x, y in candidates:
        if 0 <= x < image.width and 0 <= y < image.height:
            pixel = image.getpixel((x, y))
            if isinstance(pixel, int):
                samples.append((pixel, pixel, pixel))
            else:
                samples.append(tuple(pixel[:3]))
    if not samples:
        return (255, 255, 255)
    samples.sort()
    return samples[len(samples) // 2]


def extract_ocr_json(config: OcrConfig) -> Path:
    rapid_ocr = None
    try:
        from paddleocr import PaddleOCR
    except ImportError:
        PaddleOCR = None  # type: ignore[assignment]

    if PaddleOCR is None:
        rapid_ocr = create_rapid_ocr()
        if rapid_ocr is None:
            raise RuntimeError(
                "No OCR engine is available. Install optional OCR dependencies "
                "with pip install -e .[ocr], or install rapidocr-onnxruntime, "
                "or provide OCR JSON manually."
            )

    ocr = None
    if PaddleOCR is not None:
        try:
            ocr = PaddleOCR(use_angle_cls=True, lang=config.lang)
        except Exception:
            rapid_ocr = create_rapid_ocr()
            if rapid_ocr is None:
                raise

    if ocr is None and rapid_ocr is None:
        raise RuntimeError(
            "No OCR engine is available. Provide OCR JSON manually."
        )

    images = collect_slide_images(config.src.resolve(), config.pattern)
    config.out.mkdir(parents=True, exist_ok=True)

    for image_path in images:
        if ocr is not None:
            try:
                raw_result = run_paddle_ocr(ocr, image_path)
                blocks = paddle_result_to_blocks(raw_result, image_path=image_path)
            except Exception:
                rapid_ocr = rapid_ocr or create_rapid_ocr()
                if rapid_ocr is None:
                    raise
                ocr = None
                raw_result = run_rapid_ocr(rapid_ocr, image_path)
                blocks = rapid_result_to_blocks(raw_result, image_path=image_path)
        else:
            raw_result = run_rapid_ocr(rapid_ocr, image_path)
            blocks = rapid_result_to_blocks(raw_result, image_path=image_path)
        (config.out / f"{image_path.stem}.json").write_text(
            json.dumps({"blocks": blocks}, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

    return config.out


def create_rapid_ocr() -> Any | None:
    try:
        from rapidocr_onnxruntime import RapidOCR
    except ImportError:
        return None
    return RapidOCR()


def run_paddle_ocr(ocr: Any, image_path: Path) -> Any:
    try:
        return ocr.ocr(str(image_path), cls=True)
    except TypeError:
        return ocr.predict(str(image_path))


def run_rapid_ocr(ocr: Any, image_path: Path) -> Any:
    result, _ = ocr(str(image_path))
    return result or []


def paddle_result_to_blocks(
    raw_result: Any,
    *,
    image_path: Path | None = None,
) -> list[dict[str, Any]]:
    blocks: list[dict[str, Any]] = []
    pages = raw_result if isinstance(raw_result, list) else [raw_result]
    if len(pages) == 1 and isinstance(pages[0], list):
        lines = pages[0]
    else:
        lines = []
        for page in pages:
            if isinstance(page, list):
                lines.extend(page)

    for line in lines:
        if not isinstance(line, (list, tuple)) or len(line) < 2:
            continue
        box = line[0]
        payload = line[1]
        if not isinstance(payload, (list, tuple)) or not payload:
            continue
        text = str(payload[0])
        confidence = float(payload[1]) if len(payload) > 1 else None
        bbox = polygon_to_xywh(box)
        blocks.append(
            {
                "text": text,
                "bbox": bbox,
                "confidence": confidence,
                "font_size": estimate_font_size_from_bbox(bbox[3]),
                "color": sample_text_color(image_path, bbox),
            }
        )
    return blocks


def rapid_result_to_blocks(
    raw_result: Any,
    *,
    image_path: Path | None = None,
) -> list[dict[str, Any]]:
    blocks: list[dict[str, Any]] = []
    for line in raw_result or []:
        if not isinstance(line, (list, tuple)) or len(line) < 2:
            continue
        box = line[0]
        text = str(line[1])
        confidence = float(line[2]) if len(line) > 2 else None
        bbox = polygon_to_xywh(box)
        blocks.append(
            {
                "text": text,
                "bbox": bbox,
                "confidence": confidence,
                "font_size": estimate_font_size_from_bbox(bbox[3]),
                "color": sample_text_color(image_path, bbox),
            }
        )
    return blocks


def estimate_font_size_from_bbox(box_height_px: float) -> float:
    return max(6.0, box_height_px * 0.75)


def sample_text_color(image_path: Path | None, bbox: list[float]) -> str:
    if image_path is None:
        return "#111827"

    try:
        from PIL import Image
    except ImportError:
        return "#111827"

    x, y, w, h = [int(round(value)) for value in bbox]
    if w <= 0 or h <= 0:
        return "#111827"

    with Image.open(image_path) as source:
        image = source.convert("RGB")

    left = max(0, x)
    top = max(0, y)
    right = min(image.width, x + w)
    bottom = min(image.height, y + h)
    if right <= left or bottom <= top:
        return "#111827"

    crop = image.crop((left, top, right, bottom))
    pixels = list(crop.get_flattened_data()) if hasattr(crop, "get_flattened_data") else list(crop.convert("RGB").getdata())
    if not pixels:
        return "#111827"

    mean = tuple(sum(channel) / len(pixels) for channel in zip(*pixels))
    luminance = 0.299 * mean[0] + 0.587 * mean[1] + 0.114 * mean[2]

    def pixel_lum(rgb: tuple[int, int, int]) -> float:
        return 0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]

    if luminance < 128:
        selected = [rgb for rgb in pixels if pixel_lum(rgb) >= luminance]
    else:
        selected = [rgb for rgb in pixels if pixel_lum(rgb) <= luminance]

    if not selected:
        selected = pixels

    avg = [int(round(sum(values) / len(selected))) for values in zip(*selected)]
    return f"#{avg[0]:02X}{avg[1]:02X}{avg[2]:02X}"


def load_spec_lines(spec_file: Path | None) -> list[str]:
    if spec_file is None:
        return []
    lines = []
    for raw in spec_file.read_text(encoding="utf-8").splitlines():
        text = raw.strip()
        if text:
            lines.append(text)
    return lines


def apply_spec_corrections(
    blocks: list[dict[str, Any]],
    spec_lines: list[str],
    threshold: float = 0.72,
) -> list[dict[str, Any]]:
    if not blocks or not spec_lines:
        return blocks

    corrected: list[dict[str, Any]] = []
    remaining = list(spec_lines)
    for block in blocks:
        text = str(block.get("text", "")).strip()
        if not text:
            corrected.append(block)
            continue

        best_match = None
        best_score = 0.0
        for candidate in remaining:
            score = SequenceMatcher(None, text.lower(), candidate.lower()).ratio()
            if score > best_score:
                best_score = score
                best_match = candidate

        if best_match is not None and best_score >= threshold:
            new_block = dict(block)
            new_block["text"] = best_match
            new_block["source"] = "spec-corrected"
            corrected.append(new_block)
            remaining.remove(best_match)
        else:
            corrected.append(block)
    return corrected


def polygon_to_xywh(points: Any) -> list[float]:
    xs = [float(point[0]) for point in points]
    ys = [float(point[1]) for point in points]
    return [min(xs), min(ys), max(xs) - min(xs), max(ys) - min(ys)]


def load_slide_blocks(ocr_dir: Path, stem: str) -> list[dict[str, Any]]:
    path = ocr_dir / f"{stem}.json"
    if not path.exists():
        return []
    data = json.loads(path.read_text(encoding="utf-8-sig"))
    if isinstance(data, list):
        return data
    if isinstance(data, dict):
        blocks = data.get("blocks", [])
        if isinstance(blocks, list):
            return blocks
    raise ValueError(f"Unsupported OCR JSON schema: {path}")


def normalize_bbox(raw_bbox: Any, width: int, height: int) -> tuple[float, float, float, float]:
    if isinstance(raw_bbox, dict):
        x = float(raw_bbox.get("x", raw_bbox.get("left", 0)))
        y = float(raw_bbox.get("y", raw_bbox.get("top", 0)))
        w = float(raw_bbox.get("w", raw_bbox.get("width", 0)))
        h = float(raw_bbox.get("h", raw_bbox.get("height", 0)))
    elif isinstance(raw_bbox, list) and len(raw_bbox) == 4:
        x, y, w, h = [float(value) for value in raw_bbox]
        if w > x and h > y and (w > width * 0.5 or h > height * 0.5):
            w -= x
            h -= y
    else:
        raise ValueError(f"Unsupported bbox value: {raw_bbox!r}")

    return clamp_box(x, y, w, h, width, height)


def clamp_box(
    x: float,
    y: float,
    w: float,
    h: float,
    width: int,
    height: int,
) -> tuple[float, float, float, float]:
    x = max(0, min(x, width))
    y = max(0, min(y, height))
    w = max(1, min(w, width - x))
    h = max(1, min(h, height - y))
    return x, y, w, h


def px_to_emu(value: float, source_extent: int, slide_extent: int) -> int:
    return int(round((value / source_extent) * slide_extent))


def estimate_font_size(box_height_px: float, slide_height_px: int) -> float:
    slide_height_points = 7.5 * 72
    return max(6, min(44, box_height_px / slide_height_px * slide_height_points * 0.72))


def parse_hex_color(value: str | None) -> RGBColor | None:
    if not value:
        return None
    match = re.fullmatch(r"#?([0-9a-fA-F]{6})", str(value).strip())
    if not match:
        return None
    raw = match.group(1)
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def parse_alignment(value: Any) -> PP_ALIGN:
    normalized = str(value or "left").lower()
    if normalized == "center":
        return PP_ALIGN.CENTER
    if normalized == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT
