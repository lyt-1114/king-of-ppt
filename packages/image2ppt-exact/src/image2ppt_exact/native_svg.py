from __future__ import annotations

import base64
import re
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .exporter import natural_path_key


SVG_NS = "http://www.w3.org/2000/svg"
XLINK_NS = "http://www.w3.org/1999/xlink"


@dataclass(frozen=True)
class SvgNativeRebuildConfig:
    src: Path
    pptx_path: Path
    pattern: str = "slide_*.svg"
    width: int | None = None
    height: int | None = None
    default_font: str = "Microsoft YaHei"
    default_color: str = "#111827"


@dataclass(frozen=True)
class SvgNativeRebuildResult:
    pptx_path: Path
    log_path: Path
    slide_count: int
    text_count: int
    shape_count: int
    picture_count: int
    line_count: int
    group_count: int
    skipped_count: int


@dataclass
class SvgContext:
    slide: Any
    prs: Presentation
    svg_path: Path
    source_w: float
    source_h: float
    default_font: str
    default_color: str
    stats: dict[str, int]
    inherited: dict[str, str]
    temp_dir: Path
    dx: float = 0
    dy: float = 0
    sx: float = 1
    sy: float = 1


def rebuild_from_svg_native(config: SvgNativeRebuildConfig) -> SvgNativeRebuildResult:
    svg_files = collect_svg_files(config.src.resolve(), config.pattern)
    source_w, source_h = resolve_canvas(svg_files[0], config.width, config.height)

    prs = Presentation()
    prs.slide_width = Inches(source_w / 144)
    prs.slide_height = Inches(source_h / 144)
    blank_layout = prs.slide_layouts[6]
    stats = {
        "text": 0,
        "shape": 0,
        "picture": 0,
        "line": 0,
        "group": 0,
        "skipped": 0,
    }

    with tempfile.TemporaryDirectory() as tmp:
        temp_dir = Path(tmp)
        for svg_path in svg_files:
            slide = prs.slides.add_slide(blank_layout)
            root = ET.parse(svg_path).getroot()
            ctx = SvgContext(
                slide=slide,
                prs=prs,
                svg_path=svg_path,
                source_w=source_w,
                source_h=source_h,
                default_font=config.default_font,
                default_color=config.default_color,
                stats=stats,
                inherited={},
                temp_dir=temp_dir,
            )
            for child in root:
                add_svg_element(child, ctx)

        config.pptx_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(config.pptx_path)

    log_path = config.pptx_path.with_suffix(".svg-native-log.md")
    log_path.write_text(
        build_svg_native_log(config, svg_files, source_w, source_h, stats),
        encoding="utf-8",
    )

    return SvgNativeRebuildResult(
        pptx_path=config.pptx_path,
        log_path=log_path,
        slide_count=len(svg_files),
        text_count=stats["text"],
        shape_count=stats["shape"],
        picture_count=stats["picture"],
        line_count=stats["line"],
        group_count=stats["group"],
        skipped_count=stats["skipped"],
    )


def collect_svg_files(src: Path, pattern: str) -> list[Path]:
    if not src.exists():
        raise FileNotFoundError(f"SVG source folder not found: {src}")
    if not src.is_dir():
        raise NotADirectoryError(f"SVG source path is not a folder: {src}")
    files = [path for path in src.glob(pattern) if path.is_file()]
    files.sort(key=natural_path_key)
    if not files:
        raise FileNotFoundError(f"No SVG files found under {src} with pattern {pattern!r}")
    return files


def resolve_canvas(svg_path: Path, width: int | None, height: int | None) -> tuple[float, float]:
    if width and height:
        return float(width), float(height)
    root = ET.parse(svg_path).getroot()
    view_box = root.get("viewBox")
    if view_box:
        parts = [float(part) for part in re.split(r"[\s,]+", view_box.strip()) if part]
        if len(parts) == 4 and parts[2] > 0 and parts[3] > 0:
            return parts[2], parts[3]
    svg_w = parse_number(root.get("width"), 1920)
    svg_h = parse_number(root.get("height"), 1080)
    return svg_w, svg_h


def add_svg_element(elem: ET.Element, ctx: SvgContext) -> None:
    tag = strip_ns(elem.tag)
    if tag in {"defs", "title", "desc", "metadata", "style"}:
        return
    if tag == "g":
        ctx.stats["group"] += 1
        child_ctx = child_context(elem, ctx)
        for child in elem:
            add_svg_element(child, child_ctx)
        return
    if tag == "rect":
        add_rect(elem, ctx)
        return
    if tag in {"circle", "ellipse"}:
        add_ellipse(elem, ctx, tag)
        return
    if tag == "line":
        add_line(elem, ctx)
        return
    if tag in {"polyline", "polygon"}:
        add_poly_shape(elem, ctx, close=(tag == "polygon"))
        return
    if tag == "path":
        add_path(elem, ctx)
        return
    if tag == "text":
        add_text(elem, ctx)
        return
    if tag == "image":
        add_image(elem, ctx)
        return
    ctx.stats["skipped"] += 1


def child_context(elem: ET.Element, ctx: SvgContext) -> SvgContext:
    dx, dy, sx, sy = parse_transform(elem.get("transform", ""))
    inherited = dict(ctx.inherited)
    inherited.update(extract_style_attrs(elem))
    return SvgContext(
        slide=ctx.slide,
        prs=ctx.prs,
        svg_path=ctx.svg_path,
        source_w=ctx.source_w,
        source_h=ctx.source_h,
        default_font=ctx.default_font,
        default_color=ctx.default_color,
        stats=ctx.stats,
        inherited=inherited,
        temp_dir=ctx.temp_dir,
        dx=ctx.dx + dx * ctx.sx,
        dy=ctx.dy + dy * ctx.sy,
        sx=ctx.sx * sx,
        sy=ctx.sy * sy,
    )


def add_rect(elem: ET.Element, ctx: SvgContext) -> None:
    x = tx(parse_number(elem.get("x"), 0), ctx)
    y = ty(parse_number(elem.get("y"), 0), ctx)
    w = tw(parse_number(elem.get("width"), 0), ctx)
    h = th(parse_number(elem.get("height"), 0), ctx)
    if w <= 0 or h <= 0:
        ctx.stats["skipped"] += 1
        return
    shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if parse_number(elem.get("rx"), 0) > 0 else MSO_SHAPE.RECTANGLE
    shape = ctx.slide.shapes.add_shape(shape_kind, emu_x(x, ctx), emu_y(y, ctx), emu_w(w, ctx), emu_h(h, ctx))
    apply_fill_and_line(shape, elem, ctx)
    ctx.stats["shape"] += 1


def add_ellipse(elem: ET.Element, ctx: SvgContext, tag: str) -> None:
    if tag == "circle":
        r = parse_number(elem.get("r"), 0)
        cx = parse_number(elem.get("cx"), 0)
        cy = parse_number(elem.get("cy"), 0)
        x, y, w, h = cx - r, cy - r, r * 2, r * 2
    else:
        rx = parse_number(elem.get("rx"), 0)
        ry = parse_number(elem.get("ry"), 0)
        cx = parse_number(elem.get("cx"), 0)
        cy = parse_number(elem.get("cy"), 0)
        x, y, w, h = cx - rx, cy - ry, rx * 2, ry * 2
    x, y, w, h = tx(x, ctx), ty(y, ctx), tw(w, ctx), th(h, ctx)
    if w <= 0 or h <= 0:
        ctx.stats["skipped"] += 1
        return
    shape = ctx.slide.shapes.add_shape(MSO_SHAPE.OVAL, emu_x(x, ctx), emu_y(y, ctx), emu_w(w, ctx), emu_h(h, ctx))
    apply_fill_and_line(shape, elem, ctx)
    ctx.stats["shape"] += 1


def add_line(elem: ET.Element, ctx: SvgContext) -> None:
    x1 = tx(parse_number(elem.get("x1"), 0), ctx)
    y1 = ty(parse_number(elem.get("y1"), 0), ctx)
    x2 = tx(parse_number(elem.get("x2"), 0), ctx)
    y2 = ty(parse_number(elem.get("y2"), 0), ctx)
    line = ctx.slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        emu_x(x1, ctx),
        emu_y(y1, ctx),
        emu_x(x2, ctx),
        emu_y(y2, ctx),
    )
    apply_line(line, elem, ctx)
    ctx.stats["line"] += 1


def add_poly_shape(elem: ET.Element, ctx: SvgContext, close: bool) -> None:
    points = parse_points(elem.get("points", ""))
    if len(points) < 2:
        ctx.stats["skipped"] += 1
        return
    add_freeform(points, close, elem, ctx)


def add_path(elem: ET.Element, ctx: SvgContext) -> None:
    points, close = parse_simple_path(elem.get("d", ""))
    if len(points) < 2:
        ctx.stats["skipped"] += 1
        return
    add_freeform(points, close, elem, ctx)


def add_freeform(points: list[tuple[float, float]], close: bool, elem: ET.Element, ctx: SvgContext) -> None:
    transformed = [(tx(x, ctx), ty(y, ctx)) for x, y in points]
    start_x, start_y = transformed[0]
    builder = ctx.slide.shapes.build_freeform(emu_x(start_x, ctx), emu_y(start_y, ctx))
    vertices = [(emu_x(x, ctx), emu_y(y, ctx)) for x, y in transformed[1:]]
    if vertices:
        builder.add_line_segments(vertices, close=close)
    shape = builder.convert_to_shape()
    apply_fill_and_line(shape, elem, ctx)
    ctx.stats["shape"] += 1


def add_text(elem: ET.Element, ctx: SvgContext) -> None:
    text = collect_text(elem)
    if not text.strip():
        ctx.stats["skipped"] += 1
        return
    x = tx(parse_number(elem.get("x"), 0), ctx)
    y = ty(parse_number(elem.get("y"), 0), ctx)
    font_size = parse_number(get_attr(elem, "font-size", ctx), 24)
    width = parse_number(elem.get("width"), max(120, len(text) * font_size * 0.72))
    height = parse_number(elem.get("height"), font_size * 1.45)
    box_y = y - font_size * 0.9
    shape = ctx.slide.shapes.add_textbox(
        emu_x(x, ctx),
        emu_y(box_y, ctx),
        emu_w(tw(width, ctx), ctx),
        emu_h(th(height, ctx), ctx),
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = parse_alignment(get_attr(elem, "text-anchor", ctx))
    run = paragraph.add_run()
    run.text = text
    run.font.name = get_attr(elem, "font-family", ctx) or ctx.default_font
    run.font.size = Pt(font_size * 0.75)
    run.font.bold = str(get_attr(elem, "font-weight", ctx) or "").lower() in {"bold", "700", "800", "900"}
    run.font.italic = str(get_attr(elem, "font-style", ctx) or "").lower() == "italic"
    run.font.color.rgb = parse_color(get_attr(elem, "fill", ctx) or ctx.default_color)
    ctx.stats["text"] += 1


def add_image(elem: ET.Element, ctx: SvgContext) -> None:
    href = elem.get("href") or elem.get(f"{{{XLINK_NS}}}href")
    if not href:
        ctx.stats["skipped"] += 1
        return
    image_path = resolve_image(href, ctx)
    if image_path is None:
        ctx.stats["skipped"] += 1
        return
    x = tx(parse_number(elem.get("x"), 0), ctx)
    y = ty(parse_number(elem.get("y"), 0), ctx)
    w = tw(parse_number(elem.get("width"), ctx.source_w), ctx)
    h = th(parse_number(elem.get("height"), ctx.source_h), ctx)
    ctx.slide.shapes.add_picture(str(image_path), emu_x(x, ctx), emu_y(y, ctx), width=emu_w(w, ctx), height=emu_h(h, ctx))
    ctx.stats["picture"] += 1


def apply_fill_and_line(shape: Any, elem: ET.Element, ctx: SvgContext) -> None:
    fill = get_attr(elem, "fill", ctx)
    if fill and fill != "none":
        color = parse_color(fill)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    else:
        shape.fill.background()
    apply_line(shape, elem, ctx)


def apply_line(shape: Any, elem: ET.Element, ctx: SvgContext) -> None:
    stroke = get_attr(elem, "stroke", ctx)
    if stroke and stroke != "none":
        shape.line.color.rgb = parse_color(stroke)
        shape.line.width = Pt(parse_number(get_attr(elem, "stroke-width", ctx), 1))
    else:
        shape.line.fill.background()


def resolve_image(href: str, ctx: SvgContext) -> Path | None:
    if href.startswith("data:image/"):
        match = re.match(r"data:image/([a-zA-Z0-9.+-]+);base64,(.+)", href, re.DOTALL)
        if not match:
            return None
        suffix = ".jpg" if match.group(1).lower() in {"jpeg", "jpg"} else ".png"
        image_path = ctx.temp_dir / f"svg_image_{ctx.stats['picture'] + 1}{suffix}"
        image_path.write_bytes(base64.b64decode(match.group(2)))
        return image_path
    path = Path(href)
    if not path.is_absolute():
        path = ctx.svg_path.parent / path
    return path.resolve() if path.exists() else None


def collect_text(elem: ET.Element) -> str:
    pieces: list[str] = []
    if elem.text:
        pieces.append(elem.text)
    for child in elem:
        if strip_ns(child.tag) == "tspan":
            pieces.append(collect_text(child))
        if child.tail:
            pieces.append(child.tail)
    return "".join(pieces)


def get_attr(elem: ET.Element, name: str, ctx: SvgContext) -> str | None:
    direct = elem.get(name)
    if direct is not None:
        return direct
    styles = parse_style(elem.get("style", ""))
    if name in styles:
        return styles[name]
    return ctx.inherited.get(name)


def extract_style_attrs(elem: ET.Element) -> dict[str, str]:
    attrs = parse_style(elem.get("style", ""))
    for name in [
        "fill",
        "stroke",
        "stroke-width",
        "font-size",
        "font-family",
        "font-weight",
        "font-style",
        "text-anchor",
    ]:
        if elem.get(name) is not None:
            attrs[name] = str(elem.get(name))
    return attrs


def parse_style(style: str) -> dict[str, str]:
    result: dict[str, str] = {}
    for chunk in style.split(";"):
        if ":" not in chunk:
            continue
        key, value = chunk.split(":", 1)
        result[key.strip()] = value.strip()
    return result


def parse_transform(transform: str) -> tuple[float, float, float, float]:
    dx, dy = 0.0, 0.0
    sx, sy = 1.0, 1.0
    match = re.search(r"translate\(\s*([-\d.]+)(?:[\s,]+([-\d.]+))?\s*\)", transform)
    if match:
        dx = float(match.group(1))
        dy = float(match.group(2) or 0)
    match = re.search(r"scale\(\s*([-\d.]+)(?:[\s,]+([-\d.]+))?\s*\)", transform)
    if match:
        sx = float(match.group(1))
        sy = float(match.group(2) or sx)
    return dx, dy, sx, sy


def parse_points(value: str) -> list[tuple[float, float]]:
    nums = [float(v) for v in re.findall(r"-?\d+(?:\.\d+)?", value)]
    return [(nums[i], nums[i + 1]) for i in range(0, len(nums) - 1, 2)]


def parse_simple_path(value: str) -> tuple[list[tuple[float, float]], bool]:
    tokens = re.findall(r"[MLHVZmlhvz]|-?\d+(?:\.\d+)?", value)
    points: list[tuple[float, float]] = []
    idx = 0
    cmd = ""
    x = y = 0.0
    close = False
    while idx < len(tokens):
        token = tokens[idx]
        if re.fullmatch(r"[MLHVZmlhvz]", token):
            cmd = token
            idx += 1
            if cmd in {"Z", "z"}:
                close = True
            continue
        if cmd in {"M", "L"} and idx + 1 < len(tokens):
            x = float(tokens[idx])
            y = float(tokens[idx + 1])
            points.append((x, y))
            idx += 2
            if cmd == "M":
                cmd = "L"
            continue
        if cmd in {"m", "l"} and idx + 1 < len(tokens):
            x += float(tokens[idx])
            y += float(tokens[idx + 1])
            points.append((x, y))
            idx += 2
            if cmd == "m":
                cmd = "l"
            continue
        if cmd == "H":
            x = float(tokens[idx])
            points.append((x, y))
            idx += 1
            continue
        if cmd == "h":
            x += float(tokens[idx])
            points.append((x, y))
            idx += 1
            continue
        if cmd == "V":
            y = float(tokens[idx])
            points.append((x, y))
            idx += 1
            continue
        if cmd == "v":
            y += float(tokens[idx])
            points.append((x, y))
            idx += 1
            continue
        break
    return points, close


def parse_number(value: str | None, default: float = 0) -> float:
    if value is None:
        return float(default)
    match = re.search(r"-?\d+(?:\.\d+)?", str(value))
    return float(match.group(0)) if match else float(default)


def parse_color(value: str | None) -> RGBColor:
    raw = (value or "#000000").strip()
    if raw.startswith("rgb"):
        nums = [int(float(n)) for n in re.findall(r"\d+(?:\.\d+)?", raw)[:3]]
        if len(nums) == 3:
            return RGBColor(nums[0], nums[1], nums[2])
    raw = raw.lstrip("#")
    if len(raw) == 3:
        raw = "".join(ch * 2 for ch in raw)
    if not re.fullmatch(r"[0-9a-fA-F]{6}", raw):
        raw = "000000"
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def parse_alignment(value: str | None) -> PP_ALIGN:
    if value == "middle":
        return PP_ALIGN.CENTER
    if value == "end":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def strip_ns(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def tx(value: float, ctx: SvgContext) -> float:
    return ctx.dx + value * ctx.sx


def ty(value: float, ctx: SvgContext) -> float:
    return ctx.dy + value * ctx.sy


def tw(value: float, ctx: SvgContext) -> float:
    return value * ctx.sx


def th(value: float, ctx: SvgContext) -> float:
    return value * ctx.sy


def emu_x(value: float, ctx: SvgContext) -> int:
    return int(round((value / ctx.source_w) * ctx.prs.slide_width))


def emu_y(value: float, ctx: SvgContext) -> int:
    return int(round((value / ctx.source_h) * ctx.prs.slide_height))


def emu_w(value: float, ctx: SvgContext) -> int:
    return int(round((value / ctx.source_w) * ctx.prs.slide_width))


def emu_h(value: float, ctx: SvgContext) -> int:
    return int(round((value / ctx.source_h) * ctx.prs.slide_height))


def build_svg_native_log(
    config: SvgNativeRebuildConfig,
    svg_files: list[Path],
    source_w: float,
    source_h: float,
    stats: dict[str, int],
) -> str:
    return "\n".join(
        [
            "# SVG Native Rebuild Log",
            "",
            "- Route: structured SVG -> native PowerPoint objects.",
            "- Purpose: convert semantic SVG elements into editable PPT text boxes, shapes, lines, and pictures.",
            "- Boundary: a full-slide bitmap embedded inside SVG remains a picture, not recovered native objects.",
            "",
            f"- Source SVG folder: `{config.src.resolve()}`",
            f"- Pattern: `{config.pattern}`",
            f"- Output PPTX: `{config.pptx_path.resolve()}`",
            f"- Canvas: `{source_w:g} x {source_h:g}`",
            f"- Slides: `{len(svg_files)}`",
            f"- Native text objects: `{stats['text']}`",
            f"- Native shape objects: `{stats['shape']}`",
            f"- Picture objects: `{stats['picture']}`",
            f"- Line objects: `{stats['line']}`",
            f"- SVG group containers visited: `{stats['group']}`",
            f"- Skipped unsupported/empty elements: `{stats['skipped']}`",
            "",
        ]
    )
