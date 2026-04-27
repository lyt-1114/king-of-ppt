#!/usr/bin/env python3
"""Lightweight audit for Ultimate PPT output folders.

Checks basic output presence, source visibility, text density, and common PPTX
layout problems such as overlapping text, duplicate footers, and off-slide text.
"""

from pathlib import Path
import argparse
import re
import sys

EMU_PER_INCH = 914400


def read_pptx_text(path):
    try:
        from pptx import Presentation
    except Exception:
        return None, "python-pptx is not installed"
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        slides.append("\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text")))
    return slides, None


def content_units(text):
    """Approximate readable content length across Latin and CJK text."""
    cjk = re.findall(r"[\u4e00-\u9fff]", text)
    words = re.findall(r"[A-Za-z0-9_]+", text)
    return len(cjk) + len(words)


def shape_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text_frame.text.strip()


def shape_font_sizes(shape):
    if not getattr(shape, "has_text_frame", False):
        return []
    sizes = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip() and run.font.size is not None:
                sizes.append(run.font.size.pt)
    return sizes


def intersection_ratio(a, b):
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    ix = max(0, min(ax2, bx2) - max(ax1, bx1))
    iy = max(0, min(ay2, by2) - max(ay1, by1))
    inter = ix * iy
    if inter == 0:
        return 0
    area_a = max(1, (ax2 - ax1) * (ay2 - ay1))
    area_b = max(1, (bx2 - bx1) * (by2 - by1))
    return inter / min(area_a, area_b)


def audit_pptx_layout(pptx):
    try:
        from pptx import Presentation
    except Exception:
        return [], ["python-pptx is not installed; skipped PPTX layout audit."]

    prs = Presentation(pptx)
    errors = []
    warnings = []
    slide_w, slide_h = prs.slide_width, prs.slide_height
    margin = int(0.02 * EMU_PER_INCH)

    for slide_idx, slide in enumerate(prs.slides, 1):
        text_shapes = []
        slide_text = []
        for shape in slide.shapes:
            txt = shape_text(shape)
            if not txt:
                continue
            slide_text.append(txt)
            bounds = (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)
            text_shapes.append((bounds, txt))
            if shape.left < -margin or shape.top < -margin or shape.left + shape.width > slide_w + margin or shape.top + shape.height > slide_h + margin:
                errors.append(f"{pptx.name} slide {slide_idx}: text shape extends outside slide bounds: {txt[:40]!r}.")

            sizes = shape_font_sizes(shape)
            if sizes:
                min_size = min(sizes)
                if content_units(txt) > 4 and min_size < 8:
                    errors.append(f"{pptx.name} slide {slide_idx}: very small text ({min_size:.1f} pt): {txt[:40]!r}.")
                elif content_units(txt) > 8 and min_size < 12:
                    warnings.append(f"{pptx.name} slide {slide_idx}: small presentation text ({min_size:.1f} pt): {txt[:40]!r}.")

            area_in2 = max(0.1, (shape.width / EMU_PER_INCH) * (shape.height / EMU_PER_INCH))
            density = content_units(txt) / area_in2
            if density > 45 and content_units(txt) > 20:
                warnings.append(f"{pptx.name} slide {slide_idx}: dense text box ({density:.0f} units/in^2): {txt[:40]!r}.")

        full_units = content_units("\n".join(slide_text))
        if full_units > 150:
            warnings.append(f"{pptx.name} slide {slide_idx}: high slide text density ({full_units} content units).")

        normalized = [re.sub(r"\s+", " ", t).strip().lower() for t in slide_text]
        for repeated in ("source:", "http", "page "):
            count = sum(1 for t in normalized if repeated in t)
            if count > 1:
                warnings.append(f"{pptx.name} slide {slide_idx}: possible duplicate footer/source text ({count} matches for {repeated!r}).")

        for i, (bounds_a, text_a) in enumerate(text_shapes):
            for bounds_b, text_b in text_shapes[i + 1:]:
                ratio = intersection_ratio(bounds_a, bounds_b)
                if ratio >= 0.15:
                    errors.append(
                        f"{pptx.name} slide {slide_idx}: overlapping text boxes "
                        f"({ratio:.0%} overlap): {text_a[:28]!r} / {text_b[:28]!r}."
                    )
                elif ratio >= 0.05:
                    warnings.append(
                        f"{pptx.name} slide {slide_idx}: possible text overlap "
                        f"({ratio:.0%} overlap): {text_a[:28]!r} / {text_b[:28]!r}."
                    )

    return errors, warnings


def read_markdown_files(folder):
    docs = {}
    for md in Path(folder).glob("*.md"):
        docs[md.name.lower()] = md.read_text(encoding="utf-8", errors="ignore")
    return docs


def has_doc_or_runlog(docs, filenames, required_terms):
    for name in filenames:
        if name in docs:
            return True
    runlog = "\n".join(docs.get(name, "") for name in ("run-log.md", "skill-run.md"))
    return all(term.lower() in runlog.lower() for term in required_terms)


def audit_image2_evidence(folder):
    docs = read_markdown_files(folder)
    warnings = []

    if not docs:
        return ["Image2 audit: no markdown evidence files found."]

    if not has_doc_or_runlog(docs, ("image2-brief.md",), ("image2", "must", "editable")):
        warnings.append("Image2 audit: missing image2-brief.md or equivalent run-log section.")
    if not has_doc_or_runlog(docs, ("visual-grammar.md",), ("visual grammar", "palette", "typography")):
        warnings.append("Image2 audit: missing visual-grammar.md or equivalent run-log section.")
    if "strategy-lock.md" not in docs and "strategy lock" not in "\n".join(docs.values()).lower():
        warnings.append("Image2 audit: missing strategy-lock.md or equivalent strategy section.")
    if "execution-lock.md" not in docs and "execution lock" not in "\n".join(docs.values()).lower():
        warnings.append("Image2 audit: missing execution-lock.md or equivalent execution section.")

    all_text = "\n".join(docs.values()).lower()
    for term in ("faithful", "inspired", "upgraded", "hybrid"):
        if term in all_text:
            break
    else:
        warnings.append("Image2 audit: transfer level not found (faithful, inspired, upgraded, or hybrid).")

    for term in ("editable", "bitmap"):
        if term not in all_text:
            warnings.append(f"Image2 audit: {term} layer not documented.")

    return warnings


def audit_folder(folder, require_image2=False):
    folder = Path(folder)
    errors = []
    warnings = []
    pptx_files = list(folder.glob("*.pptx"))
    html_files = list(folder.glob("*.html"))

    if not pptx_files and not html_files:
        errors.append("No PPTX or HTML deck found.")

    for pptx in pptx_files:
        slides, err = read_pptx_text(pptx)
        if err:
            warnings.append(f"{pptx.name}: {err}; skipped PPTX text audit.")
            continue
        if len(slides) < 5:
            warnings.append(f"{pptx.name}: only {len(slides)} slides.")
        first = slides[0] if slides else ""
        if "github.com" not in first and "http" not in first:
            warnings.append(f"{pptx.name}: first slide has no visible URL/source footer.")
        for i, text in enumerate(slides, 1):
            words = re.findall(r"[\w\u4e00-\u9fff]+", text)
            if len(words) > 120:
                warnings.append(f"{pptx.name} slide {i}: high text density ({len(words)} tokens).")
        layout_errors, layout_warnings = audit_pptx_layout(pptx)
        errors.extend(layout_errors)
        warnings.extend(layout_warnings)

    for html in html_files:
        text = html.read_text(encoding="utf-8", errors="ignore")
        if "<section" not in text and "class=\"slide" not in text and "class='slide" not in text:
            warnings.append(f"{html.name}: no obvious slide sections.")
        if "100vh" not in text and "100dvh" not in text:
            warnings.append(f"{html.name}: no 100vh/100dvh viewport constraint found.")

    if not (folder / "run-log.md").exists() and not (folder / "skill-run.md").exists():
        warnings.append("No run-log.md or skill-run.md found.")

    if require_image2:
        warnings.extend(audit_image2_evidence(folder))

    return errors, warnings


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("path", help="Output folder to audit")
    parser.add_argument("--image2", action="store_true", help="Also check Image2PPT evidence files and transfer notes")
    args = parser.parse_args()
    errors, warnings = audit_folder(args.path, require_image2=args.image2)
    if errors:
        print("FAIL")
        for item in errors:
            print(f"ERROR: {item}")
        for item in warnings:
            print(f"WARN: {item}")
        return 1
    print("PASS")
    for item in warnings:
        print(f"WARN: {item}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
